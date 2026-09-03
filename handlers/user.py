# -*- coding: utf-8 -*-
import base64
import json
import logging
import re
import zipfile
import io
import os

from aiogram import Router, F
from aiogram.exceptions import TelegramBadRequest
from aiogram.filters import CommandStart, Command
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.types import Message, CallbackQuery, BufferedInputFile
from pptx import Presentation

from database import db
import texts
import keyboards as kb
import pricing
import platega
import forgetapi_client
import voice
from config import cfg

router = Router(name="user")
log = logging.getLogger(__name__)


class Topup(StatesGroup):
    waiting_amount = State()


class Promo(StatesGroup):
    waiting_code = State()


class TicketFlow(StatesGroup):
    waiting_text = State()
    waiting_reply = State()


class ChatFlow(StatesGroup):
    in_chat = State()


async def _exit_active_chat(user_id: int, state: FSMContext) -> None:
    """Выход из диалога с моделью (по /start, кнопке "В главное меню" и т.п.) —
    заодно чистит историю диалога в БД, чтобы старый контекст не тянулся (и не
    тратил токены) в следующий раз, когда пользователь зайдёт в чат с этой же
    моделью заново. Если пользователь не был в чате — просто сбрасывает состояние
    как раньше, без лишних действий."""
    if await state.get_state() == ChatFlow.in_chat.state:
        data = await state.get_data()
        model_key = data.get("model_key")
        if model_key:
            await db.clear_dialog(user_id, model_key)
    await state.clear()


# ---------------- /start ----------------

@router.message(CommandStart())
async def cmd_start(message: Message, state: FSMContext):
    await _exit_active_chat(message.from_user.id, state)
    ref_by = None
    args = message.text.split(maxsplit=1)
    if len(args) > 1 and args[1].startswith("ref"):
        try:
            candidate = int(args[1].replace("ref", ""))
            if candidate != message.from_user.id:
                ref_by = candidate
        except ValueError:
            pass
    is_new_user = await db.get_user(message.from_user.id) is None
    await db.get_or_create_user(message.from_user.id, message.from_user.username, ref_by)
    welcome_text = texts.WELCOME
    if is_new_user and cfg.signup_bonus_tokens > 0:
        bonus_str = f"{cfg.signup_bonus_tokens:,}".replace(",", " ")
        welcome_text += "\n\n" + texts.SIGNUP_BONUS_TMPL.format(n=bonus_str)
    await message.answer(welcome_text, reply_markup=kb.main_reply_kb(is_admin=message.from_user.id in cfg.admin_ids))


@router.callback_query(F.data == "menu:main")
async def cb_main(call: CallbackQuery, state: FSMContext):
    await _exit_active_chat(call.from_user.id, state)
    await call.message.edit_text(texts.WELCOME)
    await call.answer()


@router.message(Command("new_chat"))
async def cmd_new_chat(message: Message, state: FSMContext):
    data = await state.get_data()
    model_key = data.get("model_key")
    if await state.get_state() != ChatFlow.in_chat.state or not model_key:
        await message.answer("💬 Сейчас нет активного диалога с моделью. Сначала выбери модель — «🤖 Модели».")
        return
    m = pricing.MODELS.get(model_key)
    if m is None:
        await state.clear()
        await message.answer(
            "⚠️ Эта модель больше недоступна в боте. Выбери другую в «🤖 Модели».",
            reply_markup=kb.main_reply_kb(is_admin=message.from_user.id in cfg.admin_ids),
        )
        return
    await db.clear_dialog(message.from_user.id, model_key)
    await message.answer(texts.CHAT_CLEARED.format(title=m.title))


# ---------------- persistent bottom menu (reply-keyboard) ----------------
# Эти кнопки закреплены снизу экрана (как в Telegram-клиенте), а не под конкретным
# сообщением. Обработчики без явного состояния матчатся в любом состоянии FSM —
# поэтому нажатие любой из них всегда прерывает текущий сценарий (ввод суммы,
# промокода, диалог с моделью и т.п.) и ведёт в соответствующий раздел.

async def _cabinet_text(user_id: int, username: str | None) -> str:
    user = await db.get_or_create_user(user_id, username)
    lines = []
    for key, m in pricing.MODELS.items():
        remaining = await db.get_wallet(user_id, key)
        lines.append(f"• {m.title}: {remaining:.0f} токенов")
    return texts.CABINET_TMPL.format(
        user_id=user["user_id"],
        username=user["username"] or "—",
        balance=user["balance_rub"],
        min_topup=cfg.min_topup_rub,
        free_requests=user["free_requests"],
        wallets="\n".join(lines),
    )


@router.message(F.text == kb.MAIN_BTN_CABINET)
async def on_btn_cabinet(message: Message, state: FSMContext):
    await state.clear()
    text = await _cabinet_text(message.from_user.id, message.from_user.username)
    await message.answer(text, reply_markup=kb.cabinet_menu())


@router.message(F.text == kb.MAIN_BTN_MODELS)
async def on_btn_models(message: Message, state: FSMContext):
    await state.clear()
    await message.answer(texts.MODELS_CATEGORY_MENU, reply_markup=kb.categories_kb())


@router.message(F.text == kb.MAIN_BTN_SUPPORT)
async def on_btn_support(message: Message, state: FSMContext):
    await state.clear()
    await message.answer(texts.SUPPORT_MENU, reply_markup=kb.support_menu())


@router.message(F.text == kb.MAIN_BTN_REFERRAL)
async def on_btn_referral(message: Message, state: FSMContext):
    await state.clear()
    bot_user = await message.bot.get_me()
    link = f"https://t.me/{bot_user.username}?start=ref{message.from_user.id}"
    cur = await db.db().execute("SELECT COUNT(*) FROM users WHERE ref_by=?", (message.from_user.id,))
    row = await cur.fetchone()
    count = row[0] if row else 0
    text = texts.REFERRAL_TMPL.format(n=cfg.referral_free_requests, percent=cfg.referral_commission_percent, link=link, count=count)
    await message.answer(text)


@router.message(F.text == kb.MAIN_BTN_CHANNEL)
async def on_btn_channel(message: Message, state: FSMContext):
    await state.clear()
    await message.answer(texts.CHANNEL_TMPL.format(url=cfg.channel_url))


MODE_TITLES = {"normal": "💬 Обычный", "economy": "⚡ Экономный"}


@router.callback_query(F.data == "menu:mode")
async def cb_mode_menu(call: CallbackQuery, state: FSMContext):
    await state.clear()
    mode = await db.get_chat_mode(call.from_user.id)
    await call.message.edit_text(
        texts.MODE_MENU_TMPL.format(mode_title=MODE_TITLES.get(mode, mode)),
        reply_markup=kb.mode_kb(mode),
    )
    await call.answer()


@router.callback_query(F.data.startswith("mode:set:"))
async def cb_mode_set(call: CallbackQuery):
    mode = call.data.split(":")[2]
    if mode not in ("normal", "economy"):
        await call.answer()
        return
    await db.set_chat_mode(call.from_user.id, mode)
    await call.message.edit_text(
        texts.MODE_MENU_TMPL.format(mode_title=MODE_TITLES.get(mode, mode)),
        reply_markup=kb.mode_kb(mode),
    )
    await call.answer(f"✅ Режим переключён: {MODE_TITLES.get(mode, mode)}", show_alert=False)


@router.message(F.text == kb.MAIN_BTN_ADMIN)
async def on_btn_admin(message: Message, state: FSMContext):
    if message.from_user.id not in cfg.admin_ids:
        return
    await state.clear()
    await message.answer(texts.ADMIN_MENU, reply_markup=kb.admin_menu())


# ---------------- cabinet (профиль + баланс, объединено в один экран) ----------------

@router.callback_query(F.data == "menu:cabinet")
async def cb_cabinet(call: CallbackQuery):
    text = await _cabinet_text(call.from_user.id, call.from_user.username)
    await call.message.edit_text(text, reply_markup=kb.cabinet_menu())
    await call.answer()


@router.callback_query(F.data == "balance:topup")
async def cb_topup_start(call: CallbackQuery, state: FSMContext):
    await state.set_state(Topup.waiting_amount)
    await call.message.edit_text(
        texts.TOPUP_ENTER_AMOUNT.format(min_topup=cfg.min_topup_rub), reply_markup=kb.back_to_main()
    )
    await call.answer()


@router.message(Topup.waiting_amount)
async def on_topup_amount(message: Message, state: FSMContext):
    raw = message.text.replace(",", ".").strip()
    try:
        amount = float(raw)
    except ValueError:
        await message.answer(texts.TOPUP_NOT_A_NUMBER)
        return
    if amount < cfg.min_topup_rub:
        await message.answer(texts.TOPUP_TOO_SMALL.format(min_topup=cfg.min_topup_rub))
        return

    await state.clear()
    method = cfg.platega_active_methods[0] if cfg.platega_active_methods else 2
    tx_id = await db.create_transaction(message.from_user.id, amount, method)
    try:
        url = await platega.create_payment(amount, tx_id, description=f"Пополнение баланса, польз. {message.from_user.id}")
    except Exception:
        log.exception("platega create_payment failed")
        url = ""
    if not url:
        await message.answer(texts.TOPUP_FAILED, reply_markup=kb.back_to_main())
        return
    await message.answer(texts.TOPUP_LINK_READY.format(amount=amount, url=url), reply_markup=kb.pay_link_kb(url))


@router.callback_query(F.data == "balance:history")
async def cb_balance_history(call: CallbackQuery):
    history = await db.user_transaction_history(call.from_user.id)
    if not history:
        text = "История платежей пуста."
    else:
        lines = []
        for tx in history:
            status = {"paid": "✅ оплачено", "pending": "⏳ ожидание", "failed": "❌ ошибка"}.get(tx["status"], tx["status"])
            lines.append(f"{tx['amount_rub']:.2f} ₽ — {status}")
        text = "📜 <b>Последние платежи:</b>\n\n" + "\n".join(lines)
    await call.message.edit_text(text, reply_markup=kb.cabinet_menu())
    await call.answer()


# ---------------- models: категория → модель → карточка (чат + покупка вместе) ----------------

@router.callback_query(F.data.startswith("models_cat:"))
async def cb_models_category(call: CallbackQuery):
    category = call.data.split(":", 1)[1]
    await call.message.edit_text(texts.MODELS_MENU, reply_markup=kb.models_in_category_kb(category))
    await call.answer()


@router.callback_query(F.data == "models:categories")
async def cb_models_categories(call: CallbackQuery):
    await call.message.edit_text(texts.MODELS_CATEGORY_MENU, reply_markup=kb.categories_kb())
    await call.answer()


@router.callback_query(F.data.startswith("model:"))
async def cb_model_card(call: CallbackQuery):
    model_key = call.data.split(":", 1)[1]
    m = pricing.MODELS.get(model_key)
    if m is None:
        await call.answer("Эта модель больше недоступна.", show_alert=True)
        await call.message.edit_text(texts.MODELS_CATEGORY_MENU, reply_markup=kb.categories_kb())
        return
    price_per_unit = f"{m.sell_rub_per_1m:.0f} ₽ за 1 000 000 токенов"
    if m.kind == "image":
        gen_price = pricing.package_price(model_key, m.max_tokens_per_generation)
        price_per_unit += (
            f"\nНа одну генерацию уходит до {m.max_tokens_per_generation:,} токенов "
            f"(зависит от размера картинки) — это примерно {gen_price} ₽ за 1 картинку."
        ).replace(",", " ")
    text = texts.MODEL_CARD_TMPL.format(title=m.title, description=m.description, price_per_unit=price_per_unit)
    await call.message.edit_text(text, reply_markup=kb.model_card_kb(model_key))
    await call.answer()


@router.callback_query(F.data.startswith("buy:"))
async def cb_buy_package(call: CallbackQuery):
    _, model_key, amount_s = call.data.split(":")
    m = pricing.MODELS.get(model_key)
    if m is None:
        await call.answer("Эта модель больше недоступна.", show_alert=True)
        await call.message.edit_text(texts.MODELS_CATEGORY_MENU, reply_markup=kb.categories_kb())
        return
    amount = int(amount_s)
    price = pricing.package_price(model_key, amount)
    user = await db.get_or_create_user(call.from_user.id, call.from_user.username)

    if user["balance_rub"] < price:
        missing = price - user["balance_rub"]
        await call.message.edit_text(
            texts.NOT_ENOUGH_BALANCE.format(missing=missing), reply_markup=kb.confirm_topup_kb()
        )
        await call.answer()
        return

    await db.add_balance(call.from_user.id, -price)
    await db.add_wallet(call.from_user.id, model_key, amount)
    unit = "токенов"
    await call.message.edit_text(
        texts.PACKAGE_BOUGHT_TEXT.format(amount=amount, unit=unit, title=m.title, price=price),
        reply_markup=kb.model_card_kb(model_key),
    )
    await call.answer("Готово!")

    result = await db.credit_referral_commission(call.from_user.id, price, cfg.referral_commission_percent)
    if result:
        referrer_id, commission = result
        try:
            await call.bot.send_message(
                referrer_id,
                f"💸 Твой реферал купил пакет «{m.title}» — тебе начислено {commission:.2f} ₽ "
                f"({cfg.referral_commission_percent:.0f}% от покупки) на баланс.",
            )
        except Exception:
            log.warning("failed to notify referrer %s about commission", referrer_id)


# ---------------- promo ----------------

@router.callback_query(F.data == "menu:promo")
async def cb_promo_start(call: CallbackQuery, state: FSMContext):
    await state.set_state(Promo.waiting_code)
    await call.message.edit_text(texts.PROMO_ENTER_CODE, reply_markup=kb.cabinet_menu())
    await call.answer()


@router.message(Promo.waiting_code)
async def on_promo_code(message: Message, state: FSMContext):
    await state.clear()
    ok, msg = await db.redeem_promo(message.text.strip(), message.from_user.id)
    text = texts.PROMO_APPLIED_OK.format(msg=msg) if ok else texts.PROMO_APPLIED_FAIL.format(msg=msg)
    await message.answer(text, reply_markup=kb.cabinet_menu())


# ---------------- support / tickets ----------------

@router.callback_query(F.data == "menu:support")
async def cb_support(call: CallbackQuery):
    await call.message.edit_text(texts.SUPPORT_MENU, reply_markup=kb.support_menu())
    await call.answer()


@router.callback_query(F.data == "support:agreement")
async def cb_agreement(call: CallbackQuery):
    text = cfg.user_agreement_url or texts.SUPPORT_NO_DOC
    await call.message.edit_text(f"📄 Пользовательское соглашение:\n{text}", reply_markup=kb.support_menu())
    await call.answer()


@router.callback_query(F.data == "support:privacy")
async def cb_privacy(call: CallbackQuery):
    text = cfg.privacy_policy_url or texts.SUPPORT_NO_DOC
    await call.message.edit_text(f"🔒 Политика конфиденциальности:\n{text}", reply_markup=kb.support_menu())
    await call.answer()


@router.callback_query(F.data == "support:tickets")
async def cb_tickets(call: CallbackQuery):
    tickets = await db.user_open_tickets(call.from_user.id)
    text = texts.TICKET_LIST_EMPTY if not tickets else "🎫 Ваши обращения:"
    await call.message.edit_text(text, reply_markup=kb.tickets_menu(tickets))
    await call.answer()


@router.callback_query(F.data == "ticket:new")
async def cb_ticket_new(call: CallbackQuery, state: FSMContext):
    await state.set_state(TicketFlow.waiting_text)
    await call.message.edit_text(texts.TICKET_ASK_TEXT, reply_markup=kb.back_to_main())
    await call.answer()


@router.message(TicketFlow.waiting_text)
async def on_ticket_text(message: Message, state: FSMContext):
    await state.clear()
    ticket_id = await db.create_ticket(message.from_user.id, message.text)
    await message.answer(texts.TICKET_CREATED.format(ticket_id=ticket_id), reply_markup=kb.back_to_main())
    for admin_id in cfg.admin_ids:
        try:
            await message.bot.send_message(
                admin_id,
                f"🎫 Новый тикет #{ticket_id} от {message.from_user.id} (@{message.from_user.username}):\n\n{message.text}",
            )
        except Exception:
            log.exception("failed to notify admin about new ticket")


@router.callback_query(F.data.startswith("ticket:view:"))
async def cb_ticket_view(call: CallbackQuery):
    ticket_id = int(call.data.split(":")[2])
    ticket = await db.get_ticket(ticket_id)
    if not ticket or ticket["user_id"] != call.from_user.id:
        await call.answer("Тикет не найден", show_alert=True)
        return
    msgs = await db.ticket_messages(ticket_id)
    lines = [f"{'Вы' if m['sender'] == 'user' else 'Поддержка'}: {m['text']}" for m in msgs]
    text = f"🎫 Тикет #{ticket_id} ({'открыт' if ticket['status']=='open' else 'закрыт'})\n\n" + "\n\n".join(lines)
    await call.message.edit_text(text, reply_markup=kb.ticket_view_kb(ticket_id, ticket["status"] == "open"))
    await call.answer()


@router.callback_query(F.data.startswith("ticket:reply:"))
async def cb_ticket_reply_start(call: CallbackQuery, state: FSMContext):
    ticket_id = int(call.data.split(":")[2])
    await state.set_state(TicketFlow.waiting_reply)
    await state.update_data(ticket_id=ticket_id)
    await call.message.edit_text("Введите сообщение для поддержки:", reply_markup=kb.back_to_main())
    await call.answer()


@router.message(TicketFlow.waiting_reply)
async def on_ticket_reply(message: Message, state: FSMContext):
    data = await state.get_data()
    ticket_id = data["ticket_id"]
    await state.clear()
    await db.add_ticket_message(ticket_id, "user", message.text)
    await message.answer(texts.NEW_TICKET_MSG_SAVED.format(ticket_id=ticket_id), reply_markup=kb.back_to_main())
    for admin_id in cfg.admin_ids:
        try:
            await message.bot.send_message(admin_id, f"✉️ Новое сообщение в тикете #{ticket_id}:\n\n{message.text}")
        except Exception:
            log.exception("failed to notify admin about ticket reply")


# ---------------- referral ----------------

@router.callback_query(F.data == "menu:referral")
async def cb_referral(call: CallbackQuery):
    bot_user = await call.bot.get_me()
    link = f"https://t.me/{bot_user.username}?start=ref{call.from_user.id}"
    cur = await db.db().execute("SELECT COUNT(*) FROM users WHERE ref_by=?", (call.from_user.id,))
    row = await cur.fetchone()
    count = row[0] if row else 0
    text = texts.REFERRAL_TMPL.format(n=cfg.referral_free_requests, percent=cfg.referral_commission_percent, link=link, count=count)
    await call.message.edit_text(text, reply_markup=kb.back_to_main())
    await call.answer()


# ---------------- chat with a model ----------------

MAX_DOWNLOAD_BYTES = 19 * 1024 * 1024  # запас от лимита Bot API (20 МБ)
MAX_TEXT_FILE_CHARS = 20_000  # обрезаем длинные текстовые файлы, чтобы не разорвать промпт

# Расширения/типы, которые пробуем прочитать как обычный текст
TEXT_MIME_PREFIXES = ("text/",)
TEXT_MIME_EXTRA = {
    "application/json",
    "application/xml",
    "application/x-yaml",
    "application/javascript",
}
# Расширения, которые почти всегда текстовые, даже если Telegram прислал их с
# generic mime "application/octet-stream" (так бывает с кодом/конфигами).
KNOWN_TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".log", ".ini", ".cfg", ".toml", ".env",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm", ".css", ".sql", ".sh",
    ".java", ".c", ".cpp", ".h", ".go", ".rs", ".rb", ".php", ".yaml", ".yml",
}
# Лимит на файлы совсем неизвестного бинарного формата, которые всё равно пробуем
# передать модели как есть (см. _build_document_attachment) — бот теперь принимает
# ЛЮБЫЕ файлы, но действительно огромный бинарник разумно ограничить.
MAX_BINARY_ATTACH_BYTES = 8 * 1024 * 1024


def _decode_as_text_file(data: bytes, mime_type: str, filename: str | None) -> tuple[dict, None]:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("utf-8", errors="ignore")
    truncated = len(text) > MAX_TEXT_FILE_CHARS
    text = text[:MAX_TEXT_FILE_CHARS]
    if truncated:
        text += "\n[…текст файла обрезан…]"
    return {
        "kind": "text_file",
        "mime_type": mime_type,
        "filename": filename or "file.txt",
        "text": text,
    }, None

# ZIP-архивы: не отправляем модели бинарные данные архива целиком, а распаковываем
# в памяти и склеиваем только текстовые/кодовые файлы внутри (остальное — картинки,
# бинарники и т.п. — пропускаем). Так можно скинуть, например, проект/логи/выгрузку одним архивом.
ZIP_MIME_TYPES = {"application/zip", "application/x-zip-compressed", "application/x-zip"}
ZIP_TEXT_EXTENSIONS = {
    ".txt", ".md", ".json", ".csv", ".tsv", ".xml", ".yaml", ".yml", ".ini", ".cfg", ".toml",
    ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm", ".css", ".sql", ".sh", ".log",
    ".java", ".c", ".cpp", ".h", ".go", ".rs", ".rb", ".php", ".env",
}
MAX_ZIP_ENTRY_CHARS = 4_000       # лимит на один файл внутри архива, чтобы не съесть весь бюджет одним файлом
MAX_ZIP_TOTAL_CHARS = MAX_TEXT_FILE_CHARS  # общий лимит на весь архив — как для обычного текстового файла
MAX_ZIP_ENTRIES = 30              # не разбираем совсем гигантские архивы построчно


def _build_zip_attachment(data: bytes, filename: str | None) -> tuple[dict | None, str | None]:
    """Распаковывает zip в памяти, склеивает текстовые файлы внутри (с заголовками-путями)
    в один текстовый attachment. Бинарные файлы (картинки, .exe, .pdf и т.п.) внутри архива
    пропускаются. Возвращает (attachment_dict, warning_text)."""
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile:
        return None, f"⚠️ Не получилось открыть архив «{filename or 'archive.zip'}» — похоже, файл повреждён."

    chunks: list[str] = []
    total = 0
    read_count = 0
    skipped_binary = 0
    skipped_big = 0

    for info in zf.infolist():
        if info.is_dir():
            continue
        if read_count >= MAX_ZIP_ENTRIES or total >= MAX_ZIP_TOTAL_CHARS:
            break
        name = info.filename
        ext = os.path.splitext(name)[1].lower()
        if ext not in ZIP_TEXT_EXTENSIONS:
            skipped_binary += 1
            continue
        if info.file_size > 300_000:  # не распаковываем совсем большие файлы
            skipped_big += 1
            continue
        try:
            raw = zf.read(info)
        except Exception:
            continue
        try:
            text = raw.decode("utf-8")
        except UnicodeDecodeError:
            text = raw.decode("utf-8", errors="ignore")
        truncated_entry = len(text) > MAX_ZIP_ENTRY_CHARS
        text = text[:MAX_ZIP_ENTRY_CHARS]
        remaining_budget = MAX_ZIP_TOTAL_CHARS - total
        if remaining_budget <= 0:
            break
        text = text[:remaining_budget]
        chunk = f"\n--- Файл: {name} ---\n{text}"
        if truncated_entry:
            chunk += "\n[…файл обрезан…]"
        chunks.append(chunk)
        total += len(chunk)
        read_count += 1

    if not chunks:
        return None, (
            f"⚠️ В архиве «{filename or 'archive.zip'}» не нашлось текстовых/кодовых файлов "
            "для анализа (поддерживаются txt/md/json/csv/py/js/html/css и т.п. — картинки и "
            "бинарники внутри архива пропускаются)."
        )

    combined = "".join(chunks)
    note = f"[Архив «{filename or 'archive.zip'}», распаковано файлов: {read_count}"
    if skipped_binary:
        note += f", пропущено бинарных/неподдерживаемых: {skipped_binary}"
    if skipped_big:
        note += f", пропущено слишком больших: {skipped_big}"
    note += "]\n"

    return {
        "kind": "text_file",
        "mime_type": "application/zip",
        "filename": filename or "archive.zip",
        "text": note + combined,
    }, None


async def _download_bytes(message: Message, file_id: str, file_size: int | None) -> bytes | None:
    if file_size and file_size > MAX_DOWNLOAD_BYTES:
        await message.answer(
            "⚠️ Файл слишком большой (лимит Telegram-бота — примерно 20 МБ). Пришлите файл поменьше."
        )
        return None
    try:
        file_info = await message.bot.get_file(file_id)
        buf = await message.bot.download_file(file_info.file_path)
        return buf.read()
    except Exception:
        log.exception("failed to download telegram file")
        await message.answer("⚠️ Не получилось скачать файл из Telegram. Попробуйте отправить ещё раз.")
        return None


def _build_document_attachment(data: bytes, mime_type: str, filename: str | None) -> tuple[dict | None, str | None]:
    """Возвращает (attachment_dict, warning_text). Один из них — None."""
    mime_type = mime_type or "application/octet-stream"
    ext = os.path.splitext(filename or "")[1].lower()
    is_zip_ext = ext == ".zip"
    if mime_type in ZIP_MIME_TYPES or is_zip_ext:
        return _build_zip_attachment(data, filename)

    if mime_type == "application/pdf" or ext == ".pdf":
        return {
            "kind": "document",
            "mime_type": mime_type,
            "filename": filename or "document.pdf",
            "data_b64": base64.b64encode(data).decode(),
        }, None

    if mime_type.startswith(TEXT_MIME_PREFIXES) or mime_type in TEXT_MIME_EXTRA or ext in KNOWN_TEXT_EXTENSIONS:
        return _decode_as_text_file(data, mime_type, filename)

    # docx/xlsx/pptx и любой другой zip-контейнер — вытаскиваем текст изнутри так же,
    # как из обычного .zip (Word/Excel/PowerPoint форматы — это zip с XML внутри).
    if zipfile.is_zipfile(io.BytesIO(data)):
        return _build_zip_attachment(data, filename)

    # Файл без явно текстового mime/расширения — но вдруг это всё равно текст
    # (Telegram иногда присылает code-файлы с mime application/octet-stream).
    try:
        text = data.decode("utf-8")
        return {
            "kind": "text_file",
            "mime_type": mime_type,
            "filename": filename or "file.txt",
            "text": text[:MAX_TEXT_FILE_CHARS] + ("\n[…текст файла обрезан…]" if len(text) > MAX_TEXT_FILE_CHARS else ""),
        }, None
    except UnicodeDecodeError:
        pass

    # Настоящий бинарный формат без известной структуры (картинка в экзотическом формате,
    # архив rar/7z, исполняемый файл и т.п.) — раз бот должен принимать ЛЮБЫЕ файлы, не
    # отказываем совсем, а передаём модели как есть (в base64) — часть моделей реально
    # умеет читать такие вложения через file-инструмент, часть просто не учтёт их. Ограничиваем
    # размер, чтобы не тратить бюджет промпта впустую на действительно огромный бинарник.
    if len(data) > MAX_BINARY_ATTACH_BYTES:
        return None, (
            f"⚠️ Файл «{filename or mime_type}» слишком большой для передачи модели как есть "
            f"(лимит {MAX_BINARY_ATTACH_BYTES // 1024 // 1024} МБ для нераспознанных форматов). "
            "Если это архив/документ — попробуй прислать его как .zip."
        )
    return {
        "kind": "document",
        "mime_type": mime_type,
        "filename": filename or "file",
        "data_b64": base64.b64encode(data).decode(),
    }, None


# ---------------- отправка ответа модели в виде zip-архива ----------------
# Если пользователь сам просит прислать ответ «в zip», «зип-файлом» и т.п. —
# упаковываем ответ модели в архив вместо обычного текстового сообщения.

ZIP_REQUEST_RE = re.compile(r"\bzip\b|\bзип\b", re.IGNORECASE)

CODE_BLOCK_RE = re.compile(r"```(\w+)?[ \t]*\r?\n(.*?)```", re.DOTALL)

CODE_LANG_TO_EXT = {
    "python": "py", "py": "py",
    "javascript": "js", "js": "js", "node": "js",
    "typescript": "ts", "ts": "ts",
    "jsx": "jsx", "tsx": "tsx",
    "html": "html", "htm": "html",
    "css": "css",
    "json": "json",
    "yaml": "yaml", "yml": "yaml",
    "sql": "sql",
    "bash": "sh", "sh": "sh", "shell": "sh", "zsh": "sh",
    "java": "java",
    "c": "c",
    "cpp": "cpp", "c++": "cpp",
    "go": "go", "golang": "go",
    "rust": "rs", "rs": "rs",
    "ruby": "rb", "rb": "rb",
    "php": "php",
    "xml": "xml",
    "markdown": "md", "md": "md",
    "csv": "csv",
    "toml": "toml",
    "ini": "ini",
    "env": "env",
}


def _wants_zip_response(text: str) -> bool:
    """Просит ли пользователь прислать ответ архивом (zip/зип)."""
    return bool(ZIP_REQUEST_RE.search(text or ""))


# ---------------- отправка ответа модели ЛЮБЫМ файлом (универсальный протокол) ----------------
# Модели указываем (см. FILE_TAG_INSTRUCTION): если нужно прислать файл — обернуть его
# содержимое в [FILE:имя.расширение] ... [/FILE]. Бот вылавливает такие блоки в любом
# ответе (независимо от того, просил пользователь zip или презентацию) и отправляет их
# как настоящие Telegram-документы, а не как текст в чат — так бот умеет "отправлять
# любые файлы", а не только заранее зашитые zip/pptx.

FILE_TAG_RE = re.compile(r"\[FILE:([^\]\r\n]{1,120})\]\r?\n?(.*?)\[/FILE\]", re.DOTALL)

FILE_INTENT_RE = re.compile(
    r"файл(ом|ы|е|а)?\b|документ(ом)?\b|скачат|прикреп|\.csv\b|\.docx\b|\.xlsx\b|\.txt\b|"
    r"\.json\b|export|экспорт|сохрани как|в виде файла",
    re.IGNORECASE,
)

# Отдельно ловим запросы "напиши код"/"сделай бота" и т.п. — по умолчанию модели любят
# отвечать кодом прямо в чат markdown-блоками, а хочется, чтобы такой ответ (особенно
# многофайловый проект вроде тг-бота: main.py, handlers/, requirements.txt и т.п.) сразу
# прилетал настоящими файлами, без явной просьбы "пришли файлом".
CODE_PROJECT_RE = re.compile(
    r"напиши\s+(код|скрипт|программ|бота|бот)|сделай\s+(бота|бот|скрипт|программ)|"
    r"код\s+(для|бота|телеграм|тг)|тг[\s-]?бот|телеграм[\s-]?бот|telegram[\s-]?bot|"
    r"\bmain\.py\b|requirements\.txt|структур[ауы]\s+проект|проект\s+бота|"
    r"несколько\s+файл|разбей\s+на\s+файл|handlers?\b|\bскрипт\b",
    re.IGNORECASE,
)

FILE_TAG_INSTRUCTION = (
    "\n\nЕсли по смыслу запроса нужно прислать файл (документ, код, таблицу, экспорт данных "
    "и т.п.), а не просто ответить текстом в чат — оберни СОДЕРЖИМОЕ каждого такого файла "
    "в теги [FILE:имя_файла.расширение] и [/FILE] (расширение выбери по смыслу: .txt/.csv/"
    ".md/.py/.json и т.п.). Если это программа/проект из нескольких файлов (например "
    "Telegram-бот) — не сваливай всё в один файл: сделай нормальную структуру (main.py, "
    "handlers/..., requirements.txt, config.py и т.п., по смыслу задачи) и оберни КАЖДЫЙ "
    "файл в свои теги [FILE:относительный/путь/имя.py][/FILE] — можно вернуть сколько угодно "
    "файлов подряд. Вне тегов оставь только короткий человеческий комментарий (не дублируй "
    "в нём содержимое файлов)."
)


def _wants_file_response(text: str) -> bool:
    """Похоже ли, что пользователь хочет получить файл (или код/проект), а не просто
    текстовый ответ в чат."""
    t = text or ""
    return bool(FILE_INTENT_RE.search(t) or CODE_PROJECT_RE.search(t))


def _extract_file_blocks(answer: str) -> tuple[list[tuple[str, str]], str]:
    """Достаёт все [FILE:имя]...[/FILE] блоки из ответа модели.
    Возвращает (список (имя_файла, содержимое), оставшийся_текст_без_блоков)."""
    files = [(m.group(1).strip(), m.group(2)) for m in FILE_TAG_RE.finditer(answer)]
    leftover = FILE_TAG_RE.sub("", answer).strip()
    return files, leftover


def _safe_zip_member_path(filename: str) -> str:
    """Приводит имя файла из [FILE:...] к безопасному относительному пути внутри zip
    (сохраняя папки вроде handlers/user.py, но без выхода за пределы архива)."""
    path = (filename or "").strip().replace("\\", "/").lstrip("/")
    parts = [p for p in path.split("/") if p not in ("", ".", "..")]
    return "/".join(parts) or "file.txt"


async def _send_answer_as_files(message: Message, files: list[tuple[str, str]], leftover: str) -> None:
    """Отправляет [FILE:...] блоки как настоящие Telegram-документы. Один файл — отдельным
    документом. Несколько файлов (типичный многофайловый проект: main.py, handlers/..., и т.п.)
    удобнее одним zip-архивом с сохранением структуры папок, чем россыпью файлов с подчёркиваниями
    вместо "/"."""
    if len(files) > 1:
        buf = io.BytesIO()
        try:
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
                for filename, content in files:
                    zf.writestr(_safe_zip_member_path(filename), content)
            doc = BufferedInputFile(buf.getvalue(), filename="project.zip")
            await message.answer_document(doc, caption=f"📦 Готовый проект ({len(files)} файлов)")
            if leftover:
                await _send_model_text(message, leftover)
            return
        except Exception:
            log.exception("failed to build/send multi-file project as zip, falling back to loose documents")
    for filename, content in files:
        safe_name = re.sub(r'[\\/*?:"<>|]', "_", filename)[:120] or "file.txt"
        try:
            doc = BufferedInputFile(content.encode("utf-8"), filename=safe_name)
            await message.answer_document(doc)
        except Exception:
            log.exception("failed to send [FILE:] block as document")
            await message.answer(f"⚠️ Не получилось отправить файл «{safe_name}», вот его содержимое:\n\n{content[:3500]}")
    if leftover:
        await _send_model_text(message, leftover)


def _build_answer_zip(answer: str) -> bytes:
    """Упаковывает ответ модели в zip. Если в ответе есть code-блоки (```...```),
    каждый становится отдельным файлом с расширением по языку, а остальной текст —
    файлом answer.txt. Если code-блоков нет — весь ответ кладётся в answer.txt."""
    blocks = list(CODE_BLOCK_RE.finditer(answer))
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        if blocks:
            for i, match in enumerate(blocks, start=1):
                lang = (match.group(1) or "").strip().lower()
                code = match.group(2)
                ext = CODE_LANG_TO_EXT.get(lang, "txt")
                fname = f"file_{i}.{ext}" if len(blocks) > 1 else f"file.{ext}"
                zf.writestr(fname, code)
            leftover = CODE_BLOCK_RE.sub("", answer).strip()
            if leftover:
                zf.writestr("answer.txt", leftover)
        else:
            zf.writestr("answer.txt", answer or "")
    return buf.getvalue()


async def _send_answer_as_zip(message: Message, answer: str) -> bool:
    """Пытается отправить ответ модели архивом. Возвращает True при успехе."""
    try:
        zip_bytes = _build_answer_zip(answer)
        doc = BufferedInputFile(zip_bytes, filename="response.zip")
        await message.answer_document(doc, caption="📦 Ответ модели в архиве")
        return True
    except Exception:
        log.exception("failed to build/send zip response")
        return False


# ---------------- генерация презентаций (.pptx) моделью ----------------
# Если пользователь просит сделать презентацию/слайды — просим модель ответить
# структурированным JSON (план слайдов), а не обычным текстом, и по этому JSON
# собираем настоящий .pptx-файл через python-pptx, который и отправляем.

PRESENTATION_REQUEST_RE = re.compile(
    r"презентац|слайд|pptx|powerpoint|power\s*point|повер\s*поинт", re.IGNORECASE
)

PRESENTATION_JSON_INSTRUCTION = (
    "Сформируй план презентации по запросу выше и ответь СТРОГО одним JSON-объектом "
    "без markdown-разметки (без ```), без пояснений до или после — только сам JSON. Формат:\n"
    '{"title": "заголовок презентации", "subtitle": "короткий подзаголовок (можно пустую строку)", '
    '"slides": [{"title": "заголовок слайда", "bullets": ["короткий пункт 1", "короткий пункт 2"]}]}\n'
    "Сделай от 5 до 12 слайдов. Каждый bullet — короткая мысль (не абзац), 3-6 bullets на слайд."
)


def _wants_presentation(text: str) -> bool:
    """Просит ли пользователь сделать презентацию/слайды."""
    return bool(PRESENTATION_REQUEST_RE.search(text or ""))


def _extract_json_object(text: str) -> dict | None:
    """Достаёт JSON-объект из ответа модели: снимает возможные ```-обёртки и
    берёт содержимое между первой '{' и последней '}'."""
    if not text:
        return None
    stripped = text.strip()
    fence = re.match(r"^```(?:json)?\s*\n?(.*?)\n?```$", stripped, re.DOTALL)
    if fence:
        stripped = fence.group(1).strip()
    start = stripped.find("{")
    end = stripped.rfind("}")
    if start == -1 or end == -1 or end <= start:
        return None
    try:
        data = json.loads(stripped[start:end + 1])
    except json.JSONDecodeError:
        return None
    return data if isinstance(data, dict) else None


def _build_pptx(data: dict) -> bytes:
    """Собирает .pptx из структуры {"title", "subtitle", "slides": [{"title","bullets"}]}."""
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = (data.get("title") or "Презентация").strip()
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = (data.get("subtitle") or "").strip()

    content_layout = prs.slide_layouts[1]  # Title and Content
    for s in data.get("slides") or []:
        if not isinstance(s, dict):
            continue
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = str(s.get("title") or "").strip()
        bullets = [str(b).strip() for b in (s.get("bullets") or []) if str(b).strip()]
        body = slide.placeholders[1].text_frame
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                body.add_paragraph().text = b
        else:
            body.text = str(s.get("text") or "").strip()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def _send_answer_as_pptx(message: Message, answer: str) -> bool:
    """Пытается распарсить ответ модели как план презентации и отправить .pptx.
    Возвращает True при успехе."""
    data = _extract_json_object(answer)
    if not data or not (data.get("slides")):
        return False
    try:
        pptx_bytes = _build_pptx(data)
    except Exception:
        log.exception("failed to build pptx from model answer")
        return False
    raw_title = (data.get("title") or "presentation").strip() or "presentation"
    safe_title = re.sub(r'[\\/*?:"<>|]', "_", raw_title)[:60]
    try:
        doc = BufferedInputFile(pptx_bytes, filename=f"{safe_title}.pptx")
        await message.answer_document(doc, caption=f"📊 {raw_title}")
    except Exception:
        log.exception("failed to send pptx document")
        return False
    return True


# ---------------- безопасная отправка текста модели в Telegram (HTML parse_mode) ----------------
# Бот на default parse_mode=HTML (см. bot.py). Если ответ модели содержит что-то похожее на
# HTML-тег (например код с `<script>`, `<your_api_key>`, математику "a < b" и т.п.) — Telegram
# отказывается парсить сообщение ("can't parse entities"), а если это исключение никак не
# ловить — сообщение "Думаю..." так и останется висеть навсегда (то самое "бот часами не
# отвечает", хотя в логах ForgetAPI запрос давно отработал). Поэтому здесь мы: 1) экранируем
# спецсимволы, 2) конвертируем базовый markdown модели (**bold**, `code`, ```блоки```,
# заголовки, списки) в настоящее форматирование Telegram вместо голых символов, 3) на случай,
# если что-то всё равно не распарсится — подстраховываемся обычным текстом без разметки.

_FENCE_RE = re.compile(r"```(\w*)\r?\n?(.*?)```", re.DOTALL)
_INLINE_CODE_RE = re.compile(r"`([^`\n]+?)`")
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*|__(.+?)__", re.DOTALL)
_ITALIC_RE = re.compile(r"(?<![\*\w])\*(?!\s)([^\*\n]+?)(?<!\s)\*(?!\*)")
_HEADER_RE = re.compile(r"^[ \t]{0,3}#{1,6}[ \t]+(.+)$", re.MULTILINE)
_BULLET_RE = re.compile(r"^[ \t]*[-*][ \t]+", re.MULTILINE)

# ---------------- LaTeX → обычный текст (страховка) ----------------
# Telegram НЕ умеет рендерить LaTeX вообще — ни \[...\], ни \(...\), ни \frac{}{}
# и т.п. Модели (особенно при решении задач/формул) любят отвечать настоящим
# LaTeX'ом, из-за чего пользователь видит вместо красивой формулы кашу из
# бэкслэшей и фигурных скобок (см. репорт: "\[m_2=0{,}5-1000\cdot1{,}5..."). Мы
# просим модель не делать так в SYSTEM_PROMPT (forgetapi_client.py), но это не
# 100% гарантия — часть моделей всё равно иногда срывается на LaTeX, особенно
# на физике/математике. Поэтому здесь, как и с HTML-сущностями выше, ставим
# отдельную страховку: конвертируем самые частые LaTeX-конструкции в читаемый
# обычный текст ПЕРЕД тем, как текст уйдёт в HTML-экранирование/разметку.
# Работает только вне ``` code ``` и `inline` блоков — там LaTeX может быть
# осознанной частью кода/примера, трогать его не нужно.

_LATEX_MATH_DELIM_RE = re.compile(r"\\[\[\]()]")  # \[ \] \( \)  — просто разделители режима формулы
_LATEX_COMMA_RE = re.compile(r"\{,\}")  # "0{,}5" -> "0,5" (экранированная десятичная запятая)
_LATEX_FRAC_RE = re.compile(r"\\d?frac\{([^{}]*)\}\{([^{}]*)\}")  # \frac{a}{b} и \dfrac{a}{b}
_LATEX_SQRT_BRACE_RE = re.compile(r"\\sqrt\{([^{}]*)\}")
_LATEX_SUPSUB_RE = re.compile(r"([\^_])\{([^{}]*)\}")  # x^{-4} -> x^-4, a_{max} -> a_max
_LATEX_WRAP_CMD_RE = re.compile(
    r"\\(text|boxed|mathrm|mathbf|mathit|mathbb|mathcal|boldsymbol|"
    r"textbf|textit|operatorname|hbox|mbox|displaystyle)\{([^{}]*)\}"
)
_LATEX_NOOP_RE = re.compile(r"\\(left|right|displaystyle|nonumber|notag)\b")
_LATEX_SPACING_RE = re.compile(r"\\[,;!]|\\quad|\\qquad")
_LATEX_FALLBACK_CMD_RE = re.compile(r"\\([a-zA-Z]+)")  # что осталось нераспознанным

_LATEX_SYMBOLS = {
    "cdot": "×", "times": "×", "div": "÷", "pm": "±", "mp": "∓",
    "leq": "≤", "le": "≤", "geq": "≥", "ge": "≥", "neq": "≠", "ne": "≠",
    "approx": "≈", "equiv": "≡", "sim": "∼", "infty": "∞", "circ": "°",
    "sum": "Σ", "int": "∫", "partial": "∂", "nabla": "∇", "sqrt": "√",
    "rightarrow": "→", "to": "→", "Rightarrow": "⇒", "leftarrow": "←",
    "in": "∈", "notin": "∉", "subset": "⊂", "cup": "∪", "cap": "∩",
    "forall": "∀", "exists": "∃", "emptyset": "∅", "propto": "∝",
    "pi": "π", "alpha": "α", "beta": "β", "gamma": "γ", "delta": "δ",
    "Delta": "Δ", "theta": "θ", "lambda": "λ", "mu": "μ", "sigma": "σ",
    "Sigma": "Σ", "omega": "ω", "Omega": "Ω", "phi": "φ", "varphi": "φ",
    "degree": "°", "%": "%",
}
_LATEX_SYMBOL_RE = re.compile(
    r"\\(" + "|".join(re.escape(k) for k in sorted(_LATEX_SYMBOLS, key=len, reverse=True)) + r")(?![a-zA-Z])"
)

# Языки, при которых код внутри ```lang ... ``` — реальный программный код (а не
# формула/расчёт), и трогать его LaTeX-страховкой нельзя: обратные слэши, `_`, `{}`
# там значимы (регулярки, f-строки, словари и т.п.). Если тег языка не входит в этот
# список (включая пустой тег — модель просто дала блок без языка под формулу, как
# просит SYSTEM_PROMPT) — считаем блок формулой/расчётом и прогоняем через _delatex,
# на случай если модель всё же оставила в нём LaTeX-мусор.
_PROGRAMMING_FENCE_LANGS = {
    "python", "py", "py3", "js", "javascript", "jsx", "ts", "typescript", "tsx",
    "java", "c", "cpp", "c++", "csharp", "cs", "go", "golang", "rust", "rs",
    "php", "ruby", "rb", "swift", "kotlin", "kt", "sql", "bash", "sh", "shell",
    "zsh", "powershell", "html", "css", "json", "yaml", "yml", "xml", "toml",
    "ini", "r", "matlab", "perl", "lua", "scala", "dart", "dockerfile", "makefile",
}


def _delatex(text: str) -> str:
    """Превращает частый LaTeX в читаемый обычный текст (см. комментарий выше)."""
    if "\\" not in text and "{,}" not in text:
        return text  # быстрый выход — в подавляющем большинстве ответов LaTeX нет вообще

    text = _LATEX_MATH_DELIM_RE.sub("", text)
    text = _LATEX_COMMA_RE.sub(",", text)
    # \frac и \sqrt со скобками гоняем в пару проходов — так подхватываются и
    # простые "соседние" вложенности вроде \frac{\sqrt{2}}{2}.
    for _ in range(3):
        text = _LATEX_FRAC_RE.sub(r"(\1)/(\2)", text)
        text = _LATEX_SQRT_BRACE_RE.sub(r"√(\1)", text)
        text = _LATEX_WRAP_CMD_RE.sub(
            lambda m: f"**{m.group(2)}**" if m.group(1) == "boxed" else m.group(2), text
        )
    text = _LATEX_SUPSUB_RE.sub(r"\1\2", text)
    text = _LATEX_SYMBOL_RE.sub(lambda m: _LATEX_SYMBOLS[m.group(1)], text)
    text = _LATEX_NOOP_RE.sub("", text)
    text = _LATEX_SPACING_RE.sub(" ", text)
    # То, что не распознали явно (\Delta, \varnothing и т.п.) — хотя бы убираем
    # бэкслэш, чтобы не висели голые "\команда" в ответе.
    text = _LATEX_FALLBACK_CMD_RE.sub(r"\1", text)
    return text


def _html_escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _format_for_telegram(text: str) -> str:
    """Готовит текст модели к отправке с parse_mode=HTML: экранирует спецсимволы и
    конвертирует базовый markdown в HTML-теги Telegram."""
    if not text:
        return text

    fences: list[str] = []

    def _stash_fence(m: re.Match) -> str:
        lang = (m.group(1) or "").strip().lower()
        content = m.group(2)
        if lang not in _PROGRAMMING_FENCE_LANGS:
            # Скорее всего формула/расчёт (или блок вообще без языка) — подчищаем
            # возможный LaTeX-мусор так же, как и в обычном тексте.
            content = _delatex(content)
        code = _html_escape(content)
        fences.append(f"<pre>{code}</pre>")
        return f"\x00FENCE{len(fences) - 1}\x00"

    text = _FENCE_RE.sub(_stash_fence, text)

    inline_codes: list[str] = []

    def _stash_inline(m: re.Match) -> str:
        code = _html_escape(m.group(1))
        inline_codes.append(f"<code>{code}</code>")
        return f"\x00CODE{len(inline_codes) - 1}\x00"

    text = _INLINE_CODE_RE.sub(_stash_inline, text)

    # LaTeX-страховка — только вне code/inline-code блоков, которые уже спрятаны выше.
    text = _delatex(text)

    # Всё, что осталось (обычный текст ответа) — экранируем перед конвертацией markdown,
    # иначе случайный "<тег>" внутри обычного текста сломает парсинг Telegram.
    text = _html_escape(text)

    text = _BOLD_RE.sub(lambda m: f"<b>{m.group(1) or m.group(2)}</b>", text)
    text = _ITALIC_RE.sub(r"<i>\1</i>", text)
    text = _HEADER_RE.sub(r"<b>\1</b>", text)
    text = _BULLET_RE.sub("• ", text)

    for i, block in enumerate(fences):
        text = text.replace(f"\x00FENCE{i}\x00", block)
    for i, block in enumerate(inline_codes):
        text = text.replace(f"\x00CODE{i}\x00", block)

    return text


# ---------------- разбивка длинных ответов на несколько сообщений ----------------
# У Telegram жёсткий лимит 4096 символов на сообщение. Раньше это никак не
# учитывалось: длинный ответ модели (например подробное решение задачи, как в
# репорте) просто падал с TelegramBadRequest на попытке отправить/отредактировать
# сообщение, fallback на обычный текст падал по той же причине (текст всё ещё
# длиннее лимита) — и пользователь оставался с вечно висящим "Думаю...", а ответ
# модели пропадал молча (в лог уходило только exception). Разбиваем заранее,
# стараясь не резать посреди ``` code ``` блока и по возможности резать на
# границах абзацев/строк, а не посреди слова.

_TG_MESSAGE_LIMIT = 4096
_TG_SAFE_CHUNK = 3500  # с запасом под HTML-теги, которые формирование добавит поверх


def _split_text_for_telegram(text: str, limit: int = _TG_SAFE_CHUNK) -> list[str]:
    """Режет сырой (ещё не HTML-форматированный) текст модели на части не длиннее
    limit символов, стараясь не разрывать ``` code ``` блоки и резать по границам
    строк/абзацев."""
    if len(text) <= limit:
        return [text]

    lines = text.split("\n")
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0
    fence_open = False

    def _flush():
        nonlocal current, current_len
        if current:
            chunks.append("\n".join(current))
        current = []
        current_len = 0

    for line in lines:
        is_fence_marker = line.lstrip().startswith("```")
        # Если строка сама по себе огромная (длиннее лимита) — режем её жёстко по
        # символам, иначе она никогда не поместится ни в один чанк.
        if len(line) > limit:
            for i in range(0, len(line), limit):
                _flush()
                chunks.append(line[i:i + limit])
            if is_fence_marker:
                fence_open = not fence_open
            continue

        extra = len(line) + (1 if current else 0)
        if current_len + extra > limit and not fence_open:
            _flush()
            extra = len(line)

        current.append(line)
        current_len += extra
        if is_fence_marker:
            fence_open = not fence_open

    _flush()

    # Финальная страховка: если внутри одного ``` code ``` блока текста оказалось
    # больше лимита (fence-aware логика выше специально его не резала, чтобы не
    # ломать разметку кода) — всё равно режем жёстко по символам. Пусть лучше один
    # блок кода придёт разбитым на несколько сообщений, чем ответ не придёт вовсе.
    hard_limit = limit * 2
    final: list[str] = []
    for chunk in chunks:
        if len(chunk) <= hard_limit:
            final.append(chunk)
        else:
            final.extend(chunk[i:i + limit] for i in range(0, len(chunk), limit))
    return final


async def _send_one_chunk(target, formatted: str, raw: str, edit: bool) -> None:
    """Отправляет один уже отформатированный чанк ответа с фоллбэком на обычный текст
    без разметки, если Telegram откажется парсить HTML (см. комментарий выше)."""
    try:
        if edit:
            await target.edit_text(formatted)
        else:
            await target.answer(formatted)
    except TelegramBadRequest:
        log.exception("failed to send HTML-formatted model answer, falling back to plain text")
        try:
            if edit:
                await target.edit_text(raw, parse_mode=None)
            else:
                await target.answer(raw, parse_mode=None)
        except Exception:
            log.exception("failed to send even plain-text fallback of model answer")


async def _send_model_text(target, text: str, edit: bool = False) -> None:
    """Отправляет текстовый ответ модели пользователю. target — либо Message (тогда
    edit должен быть False, вызывается target.answer), либо уже отправленное "Думаю..."
    сообщение (тогда edit=True, вызывается target.edit_text). Форматирует markdown
    в HTML, режет длинные ответы на несколько сообщений (лимит Telegram — 4096
    символов), а если Telegram всё равно откажется парсить результат — не роняем ответ
    молча (это и есть причина "бот висит часами"), а подстраховываемся и шлём тем же
    текстом без форматирования вообще."""
    text = text or ""
    if not text.strip():
        # Telegram отклоняет edit_text/answer с пустым текстом (Bad Request: message text
        # is empty), из-за чего сообщение "Думаю..." зависало навсегда, а модель могла
        # вернуть пустую строку (например, если весь ответ ушёл в [FILE:]-блок или API
        # отдал пустой content). Подставляем заглушку вместо падения.
        text = "⚠️ Модель вернула пустой ответ."

    raw_chunks = _split_text_for_telegram(text)
    for i, raw_chunk in enumerate(raw_chunks):
        formatted = _format_for_telegram(raw_chunk)
        if not formatted.strip():
            formatted = raw_chunk
        # Первый чанк редактирует "Думаю...", остальные (если ответ пришлось резать)
        # уходят уже новыми сообщениями — редактировать одно сообщение "в несколько
        # сообщений" невозможно.
        await _send_one_chunk(target, formatted, raw_chunk, edit=edit and i == 0)


async def _run_chat_turn(message: Message, state: FSMContext, text: str, attachments: list[dict] | None = None):
    """Общая логика одного хода чата с текстовой моделью: списание кошелька/бесплатных
    запросов, обращение к ForgetAPI, сохранение истории и ответ пользователю."""
    data = await state.get_data()
    model_key = data["model_key"]
    m = pricing.MODELS[model_key]

    remaining = await db.get_wallet(message.from_user.id, model_key)
    used_free = False
    if remaining < 1 and m.referral_eligible:
        # реферальные бесплатные запросы действуют только на GPT-5.6 Luna
        used_free = await db.use_free_request(message.from_user.id)
        if not used_free:
            await message.answer(texts.CHAT_NO_TOKENS.format(title=m.title))
            return
    elif remaining < 1:
        await message.answer(texts.CHAT_NO_TOKENS.format(title=m.title))
        return

    thinking = await message.answer(texts.CHAT_THINKING)
    wants_pptx = _wants_presentation(text)
    wants_file = _wants_file_response(text)
    chat_mode = await db.get_chat_mode(message.from_user.id)
    # Экономный режим: без истории — каждое сообщение как отдельный независимый диалог
    # (меньше токенов на запрос, но модель не помнит прошлые сообщения).
    history = [] if chat_mode == "economy" else await db.get_dialog(message.from_user.id, model_key)
    attachments_json = json.dumps(attachments, ensure_ascii=False) if attachments else None
    # для презентации просим модель вернуть структурированный JSON, а для обычного файла —
    # обернуть содержимое в [FILE:...][/FILE] — но в истории диалога сохраняем исходный текст
    # пользователя, а не служебную инструкцию
    if wants_pptx:
        api_text = f"{text}\n\n{PRESENTATION_JSON_INSTRUCTION}"
    elif wants_file:
        api_text = f"{text}{FILE_TAG_INSTRUCTION}"
    else:
        api_text = text
    history.append({"role": "user", "content": api_text, "attachments": attachments_json})
    try:
        answer, tokens_used = await forgetapi_client.chat(m, history)
    except Exception:
        log.exception("chat_completion failed (model=%s, attachments=%s)", model_key, bool(attachments))
        if attachments:
            # Часть моделей/аггрегатор не переваривают конкретное вложение (незнакомый формат,
            # слишком большой файл и т.п.) — вместо голого "ошибка, попробуйте позже" пробуем
            # ещё раз тем же сообщением, но без вложения, чтобы пользователь хотя бы получил
            # ответ по тексту, а не полный тупик.
            retry_history = history[:-1] + [{"role": "user", "content": api_text, "attachments": None}]
            try:
                answer, tokens_used = await forgetapi_client.chat(m, retry_history)
                answer = (
                    "⚠️ Не получилось прочитать прикреплённый файл этой моделью, "
                    "отвечаю только по тексту сообщения:\n\n" + answer
                )
                history = retry_history
                attachments_json = None
            except Exception:
                log.exception("chat_completion retry without attachment also failed (model=%s)", model_key)
                await thinking.edit_text(texts.CHAT_ERROR)
                return
        else:
            await thinking.edit_text(texts.CHAT_ERROR)
            return

    await db.add_message(message.from_user.id, model_key, "user", text, attachments_json)
    await db.add_message(message.from_user.id, model_key, "assistant", answer)
    if chat_mode == "economy":
        # ничего не подчищаем сверх обычного — историю мы и так не читали для запроса,
        # но подчищаем сохранённые сообщения, чтобы БД не росла бесполезно
        await db.clear_dialog(message.from_user.id, model_key)
    if not used_free:
        await db.spend_wallet(message.from_user.id, model_key, min(tokens_used, remaining))

    file_blocks, leftover_text = _extract_file_blocks(answer)
    if wants_pptx:
        await thinking.delete()
        sent = await _send_answer_as_pptx(message, answer)
        if not sent:
            await message.answer(
                "⚠️ Не получилось собрать презентацию из ответа модели. Вот ответ как есть:"
            )
            await _send_model_text(message, answer)
    elif file_blocks:
        # Модель сама обернула файл(ы) в [FILE:...] — отправляем их настоящими документами.
        await thinking.delete()
        await _send_answer_as_files(message, file_blocks, leftover_text)
    elif _wants_zip_response(text):
        await thinking.delete()
        sent = await _send_answer_as_zip(message, answer)
        if not sent:
            # если по какой-то причине не получилось собрать/отправить архив —
            # не теряем ответ, отправляем как обычный текст
            await _send_model_text(message, answer)
    else:
        await _send_model_text(thinking, answer, edit=True)


@router.callback_query(F.data.startswith("chat:clear:"))
async def cb_chat_clear(call: CallbackQuery):
    model_key = call.data.split(":")[2]
    m = pricing.MODELS.get(model_key)
    if m is None:
        await call.answer("Эта модель больше недоступна.", show_alert=True)
        return
    await db.clear_dialog(call.from_user.id, model_key)
    await call.answer(texts.CHAT_CLEARED.format(title=m.title), show_alert=True)


@router.callback_query(F.data.startswith("chat:"))
async def cb_chat_start(call: CallbackQuery, state: FSMContext):
    model_key = call.data.split(":", 1)[1]
    m = pricing.MODELS.get(model_key)
    if m is None:
        await call.answer("Эта модель больше недоступна.", show_alert=True)
        await call.message.edit_text(texts.MODELS_CATEGORY_MENU, reply_markup=kb.categories_kb())
        return
    await state.set_state(ChatFlow.in_chat)
    await state.update_data(model_key=model_key)
    if m.kind == "text":
        prompt = (
            f"💬 Диалог с {m.title} начат.\n"
            "Можно писать текст, присылать фото, голосовые или ЛЮБЫЕ файлы (текст/код/zip/PDF/"
            "docx/xlsx и т.п. — что угодно, бот разберётся).\n"
            "Попроси «пришли ответ в zip» — получишь архив, «сделай презентацию по теме ...» "
            "— получишь .pptx, а если попросить «пришли файлом»/«отправь .csv» и т.п. — модель "
            "может прислать готовый файл с ответом."
        )
    else:
        prompt = (
            f"🎨 Отправьте текстовое описание картинки для {m.title}, "
            "или пришлите фото (можно с подписью-уточнением) — сделаем изображение по образцу."
        )
    await call.answer()
    # ReplyKeyboardMarkup нельзя прикрепить через edit_text — отправляем новое сообщение,
    # оно и закрепит внизу экрана кнопки "Очистить диалог"/"В главное меню".
    await call.message.answer(prompt, reply_markup=kb.chat_reply_kb())


async def _active_model_or_reset(message: Message, state: FSMContext, model_key: str | None):
    """Достаёт ModelInfo по ключу из активной FSM-сессии. Если модель была убрана из каталога
    (например, во время обновления бота пользователь как раз сидел в чате с ней) —
    вместо KeyError аккуратно сбрасывает сессию и просит выбрать другую модель."""
    m = pricing.MODELS.get(model_key) if model_key else None
    if m is None:
        await state.clear()
        await message.answer(
            "⚠️ Эта модель больше недоступна в боте. Выбери другую в «🤖 Модели».",
            reply_markup=kb.main_reply_kb(is_admin=message.from_user.id in cfg.admin_ids),
        )
        return None
    return m


@router.message(ChatFlow.in_chat, F.text == kb.CHAT_BTN_CLEAR)
async def on_chat_clear_btn(message: Message, state: FSMContext):
    data = await state.get_data()
    model_key = data.get("model_key")
    m = await _active_model_or_reset(message, state, model_key)
    if not m:
        return
    await db.clear_dialog(message.from_user.id, model_key)
    await message.answer(texts.CHAT_CLEARED.format(title=m.title))


@router.message(ChatFlow.in_chat, F.text == kb.CHAT_BTN_BACK)
async def on_chat_back_btn(message: Message, state: FSMContext):
    await _exit_active_chat(message.from_user.id, state)
    await message.answer(texts.WELCOME, reply_markup=kb.main_reply_kb(is_admin=message.from_user.id in cfg.admin_ids))


async def _generate_image_turn(
    message: Message,
    model_key: str,
    m,
    prompt: str,
    ref_image: bytes | None = None,
    ref_mime: str = "image/jpeg",
):
    remaining = await db.get_wallet(message.from_user.id, model_key)
    if remaining < m.max_tokens_per_generation:
        await message.answer(texts.CHAT_NO_TOKENS.format(title=m.title))
        return
    thinking = await message.answer(texts.CHAT_THINKING)
    try:
        image_bytes, mime_type = await forgetapi_client.generate_image(
            m, prompt, ref_image=ref_image, ref_mime=ref_mime
        )
    except Exception:
        log.exception("generate_image failed")
        await thinking.edit_text(texts.CHAT_ERROR)
        return
    # Списываем максимальный расход токенов на генерацию (консервативно — по самому
    # дорогому доступному разрешению), т.к. точный расход конкретной генерации ForgetAPI не отдаёт.
    await db.spend_wallet(message.from_user.id, model_key, m.max_tokens_per_generation)
    await thinking.delete()
    ext = "png" if "png" in mime_type else "jpg"
    photo = BufferedInputFile(image_bytes, filename=f"generated.{ext}")
    await message.answer_photo(photo)


@router.message(ChatFlow.in_chat, F.text)
async def on_chat_text(message: Message, state: FSMContext):
    data = await state.get_data()
    model_key = data.get("model_key")
    m = await _active_model_or_reset(message, state, model_key)
    if not m:
        return

    if m.kind == "text":
        await _run_chat_turn(message, state, message.text)
    else:
        await _generate_image_turn(message, model_key, m, message.text)


@router.message(ChatFlow.in_chat, F.photo)
async def on_chat_photo(message: Message, state: FSMContext):
    data = await state.get_data()
    model_key = data.get("model_key")
    m = await _active_model_or_reset(message, state, model_key)
    if not m:
        return

    photo = message.photo[-1]
    data_bytes = await _download_bytes(message, photo.file_id, photo.file_size)
    if data_bytes is None:
        return

    if m.kind == "image":
        # Картиночная модель (nano-banana и т.п.): присланное фото — референс/описание
        # "по образцу", подпись (если есть) — доп. текстовое уточнение к нему.
        prompt = message.caption or "Сделай изображение по образцу этого фото."
        await _generate_image_turn(message, model_key, m, prompt, ref_image=data_bytes, ref_mime="image/jpeg")
        return

    if not m.supports_files:
        await message.answer(
            f"⚠️ Модель {m.title} не поддерживает фото/файлы — опишите текстом."
        )
        return

    attachment = {
        "kind": "image",
        "mime_type": "image/jpeg",
        "filename": "photo.jpg",
        "data_b64": base64.b64encode(data_bytes).decode(),
    }
    text = message.caption or "Опиши, что на этом изображении."
    await _run_chat_turn(message, state, text, [attachment])


@router.message(ChatFlow.in_chat, F.document)
async def on_chat_document(message: Message, state: FSMContext):
    data = await state.get_data()
    model_key = data.get("model_key")
    m = await _active_model_or_reset(message, state, model_key)
    if not m:
        return

    doc = message.document
    mime_type = doc.mime_type or "application/octet-stream"

    if mime_type.startswith("image/"):
        data_bytes = await _download_bytes(message, doc.file_id, doc.file_size)
        if data_bytes is None:
            return
        if m.kind == "image":
            prompt = message.caption or "Сделай изображение по образцу этого фото."
            await _generate_image_turn(message, model_key, m, prompt, ref_image=data_bytes, ref_mime=mime_type)
            return
        if not m.supports_files:
            await message.answer(f"⚠️ Модель {m.title} не поддерживает фото/файлы — опишите текстом.")
            return
        attachment = {
            "kind": "image",
            "mime_type": mime_type,
            "filename": doc.file_name or "image",
            "data_b64": base64.b64encode(data_bytes).decode(),
        }
        text = message.caption or "Опиши, что на этом изображении."
        await _run_chat_turn(message, state, text, [attachment])
        return

    if m.kind != "text" or not m.supports_files:
        await message.answer(
            f"🎨 Модель {m.title} принимает только текстовое описание — файлы не нужны, просто опишите словами."
        )
        return

    data_bytes = await _download_bytes(message, doc.file_id, doc.file_size)
    if data_bytes is None:
        return
    attachment, warning = _build_document_attachment(data_bytes, mime_type, doc.file_name)
    if warning:
        await message.answer(warning)
        return
    text = message.caption or f"Проанализируй прикреплённый файл «{doc.file_name or 'file'}»."
    await _run_chat_turn(message, state, text, [attachment])


async def _handle_voice_like(message: Message, state: FSMContext, file_id: str, file_size: int | None):
    """Общий обработчик голосовых сообщений (message.voice) и аудиофайлов (message.audio):
    скачивает файл, бесплатно распознаёт речь локальной моделью Whisper (см. voice.py)
    и дальше обрабатывает распознанный текст как обычное сообщение в чате."""
    data = await state.get_data()
    model_key = data.get("model_key")
    m = await _active_model_or_reset(message, state, model_key)
    if not m:
        return

    data_bytes = await _download_bytes(message, file_id, file_size)
    if data_bytes is None:
        return

    thinking = await message.answer("🎙️ Распознаю голосовое сообщение...")
    try:
        text = await voice.transcribe(data_bytes)
    except Exception:
        log.exception("voice transcription failed")
        await thinking.edit_text(
            "⚠️ Не получилось распознать голосовое сообщение. Попробуй ещё раз или напиши текстом."
        )
        return

    text = (text or "").strip()
    if not text:
        await thinking.edit_text("🤔 Не удалось разобрать речь в этом сообщении. Попробуй сказать чётче или напиши текстом.")
        return
    await thinking.delete()

    if m.kind == "text":
        await _run_chat_turn(message, state, text)
    else:
        await _generate_image_turn(message, model_key, m, text)


@router.message(ChatFlow.in_chat, F.voice)
async def on_chat_voice(message: Message, state: FSMContext):
    await _handle_voice_like(message, state, message.voice.file_id, message.voice.file_size)


@router.message(ChatFlow.in_chat, F.audio)
async def on_chat_audio(message: Message, state: FSMContext):
    await _handle_voice_like(message, state, message.audio.file_id, message.audio.file_size)


@router.message(ChatFlow.in_chat)
async def on_chat_unsupported(message: Message, state: FSMContext):
    await message.answer(
        "🤖 Этот тип сообщений пока не поддерживается. Отправьте текст, фото, голосовое или файл."
    )
